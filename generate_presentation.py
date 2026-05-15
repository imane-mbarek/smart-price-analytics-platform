from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

slides = [
    {
        'title': 'Smart Price Analytics Platform',
        'content': [
            'Comparateur de prix marocain avec analyse de données et notifications.',
            'Prototype intégrant scraping, backend Django, data mining et interface React.',
        ]
    },
    {
        'title': 'Objectifs du projet',
        'content': [
            'Comparer automatiquement les prix sur Jumia et Avito.',
            'Détecter les offres anormales et proposer une aide à la décision.',
            'Offrir une veille produit via panier et notifications.',
        ]
    },
    {
        'title': 'Architecture globale',
        'content': [
            'Frontend React -> Backend Django REST -> Base SQLite.',
            'Scraper interne pour récupérer les offres Jumia et Avito.',
            'Service Data Mining pour analyses, clusters et anomalies.',
        ],
        'note': 'Figure 1: architecture technique Web / API / DB / Data Mining'
    },
    {
        'title': 'Backend',
        'content': [
            'Django + Django REST Framework pour l’API.',
            'Modèles : Produit, HistoriquePrix, Panier, Notification, Recherche, SurveillancePrix.',
            'Scraping parallélisé et normalisation des prix en MAD.',
        ]
    },
    {
        'title': 'Fonctionnalités backend',
        'content': [
            'Recherche asynchrone et suivi de progression.',
            'Création de notifications pour variations de prix et stock.',
            'Export CSV et PDF des résultats de recherche.',
        ]
    },
    {
        'title': 'Module Data Mining',
        'content': [
            'Prétraitement, normalisation et encodage des données.',
            'K-Means, DBSCAN, PCA, détection d’anomalies et règles d’association.',
            'Génération d’insights pour faciliter le choix produit.',
        ],
        'note': 'Figure 2: pipeline Data Mining (prétraitement -> clustering -> anomalies -> règles)'
    },
    {
        'title': 'Frontend',
        'content': [
            'React SPA avec routes : Accueil, Historique, Panier, Notifications, Login.',
            'Hooks dédiés : useSearch, useAuth, usePanier, useNotifications.',
            'Interface utilisateur avec recherche, sélection multi-plateformes et visualisation.',
        ]
    },
    {
        'title': 'Pages principales',
        'content': [
            'Accueil : recherche, liste d’offres, export CSV/PDF.',
            'Panier : suivi des offres surveillées et variation de prix.',
            'Notifications : alertes de prix et disponibilité.',
        ]
    },
    {
        'title': 'Scénario utilisateur',
        'content': [
            '1. Connexion / inscription.',
            '2. Recherche de produit.',
            '3. Scraping des plateformes et analyse.',
            '4. Ajout au panier et suivi des notifications.',
        ],
        'note': 'Figure 3: workflow utilisateur (recherche -> analyse -> panier -> notification)'
    },
    {
        'title': 'Avantages et pistes',
        'content': [
            'Solution modulaire : backend, data mining et frontend séparés.',
            'Prototype prêt à étendre vers plus de marketplaces et données.',
            'Améliorer : base de données évolutive, dashboard graphique, détection de fraude avancée.',
        ]
    },
    {
        'title': 'Conclusion',
        'content': [
            'Outil académique et opérationnel pour l’analyse des prix en ligne.',
            'Il allie technique, data mining et expérience utilisateur.',
            'Bonne base pour une soutenance structurée et illustrée.',
        ]
    },
]

for slide_info in slides:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_info['title']
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, line in enumerate(slide_info['content']):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
    if slide_info.get('note'):
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame
        notes_text.text = slide_info['note']

prs.save('presentation_soutenance.pptx')
print('presentation_soutenance.pptx created')
